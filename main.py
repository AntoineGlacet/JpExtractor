import os
import openpyxl
from pptx import Presentation
from docx import Document
from nltk import word_tokenize
from nltk.corpus import stopwords
from langid.langid import LanguageIdentifier, model

def extract_text_from_excel(file_path):
    """Extracts text from an Excel file."""
    wb = openpyxl.load_workbook(file_path)
    text = ""

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text += " ".join(map(str, row)) + " "

    return text

def extract_text_from_powerpoint(file_path):
    """Extracts text from a PowerPoint file."""
    prs = Presentation(file_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "

    return text

def extract_text_from_word(file_path):
    """Extracts text from a Word file."""
    doc = Document(file_path)
    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + " "

    return text

def extract_japanese_words(text):
    """Extracts Japanese words from a given text."""
    identifier = LanguageIdentifier.from_modelstring(model, norm_probs=True)
    lang, _ = identifier.classify(text)

    if lang != 'ja':
        return []

    japanese_words = [word for word in word_tokenize(text) if word.isalpha()]
    stop_words = set(stopwords.words("japanese"))
    japanese_words = [word for word in japanese_words if word not in stop_words]

    return japanese_words

def process_directory(directory_path):
    """Processes all files in a directory and extracts Japanese words and their frequency."""
    japanese_words_frequency = {}

    for filename in os.listdir(directory_path):
        file_path = os.path.join(directory_path, filename)

        if filename.endswith(".xlsx"):
            text = extract_text_from_excel(file_path)
        elif filename.endswith(".pptx"):
            text = extract_text_from_powerpoint(file_path)
        elif filename.endswith(".docx"):
            text = extract_text_from_word(file_path)
        else:
            continue

        japanese_words = extract_japanese_words(text)

        for word in japanese_words:
            japanese_words_frequency[word] = japanese_words_frequency.get(word, 0) + 1

    return japanese_words_frequency

if __name__ == "__main__":
    # Replace the following path with the directory containing your files
    input_directory = "/path/to/your/directory"
    result = process_directory(input_directory)

    # Display the dictionary of Japanese words and their frequency
    print(result)
